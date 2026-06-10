"""Tests for the PowerPoint (.pptx) processor."""

from __future__ import annotations

import base64
import io
import sys

import pytest

from attachments._processors import pptx, processors

# Self-registration: pptx isn't wired into processors/__init__.py yet, so
# importing the module is what registers ".pptx" (and its option schema).
from attachments.deps import check_dep, clear_cache
from attachments.types import (
    ERROR_MISSING_DEPENDENCY,
    ERROR_PARSE,
    is_missing_dependency,
)

_requires_pptx = pytest.mark.skipif(
    not check_dep("pptx").available,
    reason="python-pptx not installed",
)


@pytest.fixture(autouse=True)
def _ensure_registered():
    """Re-register the pptx processor for every test.

    pptx.py is not yet wired into processors/__init__.py, so the conftest's
    autouse ``reset_processors()`` (which restores the built-in snapshot)
    would drop its registration after the first test.
    """
    pptx.register()


# A valid 1x1 PNG, constructed inline (no Pillow needed).
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
    "AAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


# ---------------------------------------------------------------------------
# Fixtures (decks built in-test with python-pptx)
# ---------------------------------------------------------------------------


@pytest.fixture
def two_slide_pptx():
    """Two slides, each with a title and a body placeholder."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # title and content
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Quarterly Review"
    slide1.placeholders[1].text = "Revenue grew 12 percent."
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Next Steps"
    slide2.placeholders[1].text = "Hire two engineers."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_with_image():
    """Slide 1 is text-only; slide 2 carries an embedded tiny PNG."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide1.shapes.title.text = "Cover"
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide2.shapes.add_picture(io.BytesIO(_TINY_PNG), Inches(1), Inches(1))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_with_table():
    """One untitled slide containing a 2x2 table and a text box."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(1))
    box.text_frame.text = "Scores below."
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(4), Inches(2)
    ).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "95"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@_requires_pptx
class TestPptxProcessor:
    def test_registered(self):
        assert ".pptx" in processors
        assert ".ppt" not in processors  # python-pptx cannot read legacy .ppt

    def test_text_and_kind(self, two_slide_pptx):
        result = processors[".pptx"](two_slide_pptx, filename="deck.pptx")

        assert "# Quarterly Review" in result["text"]
        assert "Revenue grew 12 percent." in result["text"]
        assert "# Next Steps" in result["text"]
        assert "Hire two engineers." in result["text"]
        assert result["meta"]["kind"] == "slides"
        assert result["meta"]["extra"]["slide_count"] == 2

    def test_artifact_structure(self, two_slide_pptx):
        result = processors[".pptx"](two_slide_pptx)
        for key in ("text", "images", "audio", "video", "meta"):
            assert key in result
        assert isinstance(result["images"], list)
        assert "error" not in result["meta"]

    def test_segments_slice_exactly(self, two_slide_pptx):
        result = processors[".pptx"](two_slide_pptx)
        text = result["text"]
        segments = result["meta"]["segments"]

        assert len(segments) == 2
        assert [s["kind"] for s in segments] == ["slide", "slide"]
        # Titles become both labels and "# <title>" heading lines.
        assert segments[0]["label"] == "Quarterly Review"
        assert segments[1]["label"] == "Next Steps"
        assert (
            text[segments[0]["start"] : segments[0]["end"]]
            == "# Quarterly Review\nRevenue grew 12 percent."
        )
        assert (
            text[segments[1]["start"] : segments[1]["end"]]
            == "# Next Steps\nHire two engineers."
        )
        # Slices reassemble into exactly the full text.
        assert text == "\n\n".join(text[s["start"] : s["end"]] for s in segments)

    def test_untitled_slide_gets_numbered_label(self, pptx_with_table):
        result = processors[".pptx"](pptx_with_table)
        segments = result["meta"]["segments"]
        assert segments[0]["label"] == "slide 1"
        assert "#" not in result["text"]  # no title, no heading line

    def test_table_extraction(self, pptx_with_table):
        result = processors[".pptx"](pptx_with_table)
        text = result["text"]
        seg = result["meta"]["segments"][0]
        slice_text = text[seg["start"] : seg["end"]]

        assert "Scores below." in slice_text
        assert "| Name | Score |" in slice_text
        assert "| Alice | 95 |" in slice_text

    def test_images_absent_by_default(self, pptx_with_image):
        result = processors[".pptx"](pptx_with_image)
        assert result["images"] == []

    def test_images_absent_when_false(self, pptx_with_image):
        result = processors[".pptx"](pptx_with_image, render_images=False)
        assert result["images"] == []
        assert "images_extracted" not in result["meta"]["extra"]

    def test_images_extracted_when_requested(self, pptx_with_image):
        result = processors[".pptx"](
            pptx_with_image, filename="deck.pptx", render_images=True
        )

        assert len(result["images"]) == 1
        img = result["images"][0]
        assert img["page"] == 2  # picture lives on the second slide
        assert img["mimetype"] == "image/png"
        assert img["bytes"] == _TINY_PNG
        assert img["name"].endswith(".png")
        assert "deck.pptx" in img["name"]
        assert result["meta"]["extra"]["images_extracted"] == 1
        # Text extraction still works alongside image extraction.
        assert "# Cover" in result["text"]

    def test_garbage_bytes_return_parse_error(self):
        result = processors[".pptx"](b"\x00\x01 definitely not a pptx")

        assert result["meta"]["error"]["code"] == ERROR_PARSE
        assert result["text"] == ""
        assert not is_missing_dependency(result)

    def test_empty_presentation(self):
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        result = processors[".pptx"](buf.getvalue())

        assert result["text"] == ""
        assert "error" not in result["meta"]
        assert result["meta"]["extra"]["slide_count"] == 0
        assert "segments" not in result["meta"]  # absent when not applicable

    def test_option_schema_resolves_to_render_images(self):
        from attachments._options import get_options, resolve_options

        schema = get_options(".pptx")
        assert [o.name for o in schema] == ["images"]
        assert schema[0].type == "bool_or_auto"

        kwargs, warnings = resolve_options(schema, {"images": "true"}, context=".pptx")
        assert kwargs == {"render_images": True}
        assert warnings == []

        kwargs, warnings = resolve_options(schema, {"render": False}, context=".pptx")
        assert kwargs == {"render_images": False}
        assert warnings == []


class TestPptxMissingDependency:
    """Always-runnable: simulates python-pptx being absent (see
    tests/test_processors/test_missing_deps.py for the masking pattern)."""

    def test_missing_dep_returns_typed_artifact(self, monkeypatch):
        monkeypatch.setitem(sys.modules, "pptx", None)
        clear_cache()
        try:
            result = processors[".pptx"](b"fake bytes", filename="deck.pptx")
        finally:
            clear_cache()

        error = result["meta"]["error"]
        assert error["code"] == ERROR_MISSING_DEPENDENCY
        assert "pip install attachments[pptx]" in error["message"]
        assert is_missing_dependency(result)
        assert result["text"] == ""
        assert result["meta"]["source"] == "deck.pptx"

    def test_missing_dep_default_filename(self, monkeypatch):
        monkeypatch.setitem(sys.modules, "pptx", None)
        clear_cache()
        try:
            result = processors[".pptx"](b"fake bytes")
        finally:
            clear_cache()

        assert is_missing_dependency(result)
        assert result["meta"]["source"] == "file.pptx"
