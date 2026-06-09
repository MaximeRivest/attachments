"""Conformance runner (VISION Epic B5).

Every artifact this library produces — happy path, error path, unpacked
archives, and the wire format served over HTTP — must validate against the
frozen JSON Schema in ``spec/artifact.schema.json``.

The suite is DYNAMIC: it iterates over the live processor registry, so
processors added later (pptx, images, ...) are picked up automatically as
long as a sample-file generator exists for their extension. Extensions
without a generator are skipped, never failed.
"""

from __future__ import annotations

import copy
import io
import json
import uuid
import zipfile
from collections.abc import Callable
from pathlib import Path

import jsonschema
import pytest

from attachments import att
from attachments.processors import processors as _registry
from attachments.server import _encode_artifact_for_wire, create_app
from attachments.types import missing_dep_artifact, normalize_artifact

SPEC_DIR = Path(__file__).resolve().parent.parent / "spec"
SCHEMA_PATH = SPEC_DIR / "artifact.schema.json"
VECTORS_PATH = SPEC_DIR / "dsl-test-vectors.json"

ARTIFACT_SCHEMA = json.loads(SCHEMA_PATH.read_text(encoding="utf-8"))
VALIDATOR = jsonschema.Draft202012Validator(ARTIFACT_SCHEMA)


# =============================================================================
# Helpers
# =============================================================================


def encode_for_wire(artifact: dict) -> dict:
    """Wire-encode a DEEP COPY of *artifact* (bytes -> bytes_b64).

    Reuses the server's encoder so the test exercises the exact transport
    transformation; the copy keeps the in-process artifact untouched.
    """
    return _encode_artifact_for_wire(copy.deepcopy(artifact))


def assert_conformant(artifact: dict) -> None:
    """Assert *artifact* (in-process form) validates against the schema."""
    wire = encode_for_wire(artifact)
    _assert_wire_conformant(wire)


def _assert_wire_conformant(wire: dict) -> None:
    """Assert an already wire-encoded artifact validates against the schema."""
    errors = sorted(VALIDATOR.iter_errors(wire), key=lambda e: list(e.absolute_path))
    if errors:
        details = "\n".join(
            f"- at {'/'.join(str(p) for p in e.absolute_path) or '<root>'}: {e.message}"
            for e in errors
        )
        pytest.fail(f"artifact does not conform to artifact.schema.json:\n{details}")


def _assert_no_none_in_meta(value: object, path: str = "meta") -> None:
    """IR rule: optional meta keys are ABSENT when not applicable, never None."""
    assert value is not None, (
        f"None at {path} — optional meta keys must be absent, never None"
    )
    if isinstance(value, dict):
        for key, item in value.items():
            _assert_no_none_in_meta(item, f"{path}.{key}")
    elif isinstance(value, list | tuple):
        for idx, item in enumerate(value):
            _assert_no_none_in_meta(item, f"{path}[{idx}]")


# =============================================================================
# Sample-file generators (extension -> bytes). Anything we cannot generate
# is skipped; extensions registered to the shared text processor fall back
# to plain UTF-8 text.
# =============================================================================


def _make_html() -> bytes:
    return (
        b"<html><head><title>Sample</title></head>"
        b"<body><h1>Title</h1><p>hello conformance</p></body></html>"
    )


def _make_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pdf() -> bytes:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _make_docx() -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph("hello conformance")
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello conformance"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_image(pil_format: str) -> Callable[[], bytes]:
    def make() -> bytes:
        from PIL import Image

        buf = io.BytesIO()
        Image.new("RGB", (4, 4), (255, 0, 0)).save(buf, format=pil_format)
        return buf.getvalue()

    return make


GENERATORS: dict[str, Callable[[], bytes]] = {
    ".html": _make_html,
    ".htm": _make_html,
    ".xlsx": _make_xlsx,
    ".pdf": _make_pdf,
    ".docx": _make_docx,
    ".pptx": _make_pptx,
    ".png": _make_image("PNG"),
    ".jpg": _make_image("JPEG"),
    ".jpeg": _make_image("JPEG"),
    ".gif": _make_image("GIF"),
    ".bmp": _make_image("BMP"),
    ".webp": _make_image("WEBP"),
    ".tiff": _make_image("TIFF"),
    ".tif": _make_image("TIFF"),
}


def _generate(ext: str) -> bytes:
    """Build tiny real file bytes for *ext*, or skip if we cannot."""
    generator = GENERATORS.get(ext)
    if generator is None:
        # Extension registered to the shared text processor: plain text works.
        return f"hello conformance via {ext}\n".encode()
    try:
        return generator()
    except Exception as exc:  # missing optional dep, template, ...
        pytest.skip(f"cannot generate a sample {ext} file: {exc}")


def _generatable_extensions() -> list[str]:
    """Extensions in the LIVE registry that we know how to generate.

    Dynamic on purpose: no hardcoded processor list, no fixed count —
    processors registered by other modules are picked up automatically.
    """
    text_proc = _registry.get("__text__")
    exts: list[str] = []
    for ext, proc in _registry.items():
        if not ext.startswith("."):
            continue  # sentinel keys like __text__
        if ext in GENERATORS or (text_proc is not None and proc is text_proc):
            exts.append(ext)
    return sorted(exts)


# =============================================================================
# Multipart / WSGI helpers (same pattern as test_service_server_integration)
# =============================================================================


def _build_multipart_body(*, filename: str, content: bytes):
    boundary = f"----attachments-{uuid.uuid4().hex}"
    chunks: list[bytes] = [
        f"--{boundary}\r\n".encode(),
        (
            f'Content-Disposition: form-data; name="file"; '
            f'filename="{filename}"\r\n'
            f"Content-Type: application/octet-stream\r\n\r\n"
        ).encode(),
        content,
        b"\r\n",
        f"--{boundary}--\r\n".encode(),
    ]
    return b"".join(chunks), f"multipart/form-data; boundary={boundary}"


def _call_wsgi(app, method: str, path: str, *, body: bytes = b"", headers=None):
    """Minimal WSGI environ builder + caller. Returns (status, json_body)."""
    environ = {
        "REQUEST_METHOD": method,
        "PATH_INFO": path,
        "CONTENT_LENGTH": str(len(body)),
        "wsgi.input": io.BytesIO(body),
    }
    for key, value in (headers or {}).items():
        env_key = "HTTP_" + key.upper().replace("-", "_")
        if key.lower() == "content-type":
            env_key = "CONTENT_TYPE"
        environ[env_key] = value

    captured: dict = {}

    def start_response(status, response_headers):
        captured["status"] = status
        captured["headers"] = response_headers

    payload = b"".join(app(environ, start_response))
    return captured["status"], json.loads(payload)


# =============================================================================
# (1) Happy path: every registered processor with a generatable extension
# =============================================================================


@pytest.mark.parametrize("ext", _generatable_extensions())
def test_processor_happy_path_artifacts_conform(ext: str, tmp_path: Path):
    data = _generate(ext)
    path = tmp_path / f"sample{ext}"
    path.write_bytes(data)

    artifacts = att(str(path))

    assert artifacts, f"att() returned no artifacts for {ext}"
    for artifact in artifacts:
        assert_conformant(artifact)
        # (7) Optional-keys rule: no None anywhere in meta on the happy path.
        _assert_no_none_in_meta(artifact["meta"])


# =============================================================================
# (2) Error artifacts conform too
# =============================================================================


class TestErrorArtifactsConform:
    def test_nonexistent_path(self):
        artifacts = att("/nonexistent/dir-conformance/missing.pdf")

        assert len(artifacts) == 1
        artifact = artifacts[0]
        assert artifact["meta"]["error"]["code"] == "unpack-error"
        assert_conformant(artifact)

    def test_pdf_processor_on_garbage_bytes(self):
        proc = _registry.get(".pdf")
        if proc is None:
            pytest.skip("no .pdf processor registered")

        artifact = proc(b"this is definitely not a PDF", filename="garbage.pdf")
        artifact = normalize_artifact(artifact, "garbage.pdf")

        assert artifact["meta"].get("error"), "expected a typed error artifact"
        assert_conformant(artifact)

    def test_missing_dependency_artifact(self):
        artifact = missing_dep_artifact("report.docx", "docx")

        assert artifact["meta"]["error"]["code"] == "missing-dependency"
        assert "pip install" in artifact["meta"]["error"]["message"]
        assert_conformant(artifact)

    def test_service_only_without_api_key(self, tmp_path: Path):
        path = tmp_path / "note.txt"
        path.write_bytes(b"hello")

        artifacts = att(str(path), prefer="service-only")

        assert len(artifacts) == 1
        artifact = artifacts[0]
        assert artifact["meta"]["error"]["code"] == "service-error"
        assert_conformant(artifact)


# =============================================================================
# (3) Unpacked archives: every artifact from a zip conforms
# =============================================================================


def test_zip_unpack_artifacts_conform(tmp_path: Path):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("hello.txt", b"hello from the zip\n")
        zf.writestr("nested/notes.md", b"# Notes\n\nbody from the zip\n")
    zip_path = tmp_path / "bundle.zip"
    zip_path.write_bytes(buf.getvalue())

    artifacts = att(str(zip_path))

    assert len(artifacts) >= 2
    for artifact in artifacts:
        assert_conformant(artifact)
        _assert_no_none_in_meta(artifact["meta"])


# =============================================================================
# (4) Wire format: the WSGI /process response body IS a conformant artifact
# =============================================================================


@pytest.mark.parametrize("ext", [".txt", ".pdf"])
def test_wsgi_process_response_conforms(ext: str):
    # ATTACHMENTS_* env (incl. server key) is scrubbed by the autouse fixture.
    app = create_app()
    body, content_type = _build_multipart_body(
        filename=f"sample{ext}", content=_generate(ext)
    )

    status, payload = _call_wsgi(
        app,
        "POST",
        "/process",
        body=body,
        headers={"Content-Type": content_type},
    )

    assert status.startswith("200")
    # The response body is already in wire form: validate it directly.
    _assert_wire_conformant(payload)
    assert payload["meta"]["source"] == f"sample{ext}"


# =============================================================================
# (5) The schema file itself is valid Draft 2020-12
# =============================================================================


def test_artifact_schema_is_valid_draft_2020_12():
    jsonschema.Draft202012Validator.check_schema(ARTIFACT_SCHEMA)


# =============================================================================
# (6) DSL test vectors: structural sanity (full run lives in test_dsl_vectors)
# =============================================================================


def test_dsl_test_vectors_parse_and_have_required_keys():
    data = json.loads(VECTORS_PATH.read_text(encoding="utf-8"))

    vectors = data["vectors"]
    assert isinstance(vectors, list)
    assert vectors, "vector file must not be empty"
    for vector in vectors:
        missing = {"input", "source", "options"} - vector.keys()
        assert not missing, (
            f"vector {vector.get('name', vector)!r} is missing keys: {missing}"
        )
