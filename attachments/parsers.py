"""File parsing logic."""

from abc import ABC, abstractmethod
from .exceptions import ParsingError
from .utils import parse_index_string # Import the new utility

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

class BaseParser(ABC):
    """Abstract base class for file parsers."""
    @abstractmethod
    def parse(self, file_path, indices=None):
        """Parses the file and returns structured content."""
        pass

class ParserRegistry:
    """Manages registration and retrieval of parsers."""
    def __init__(self):
        self.parsers = {}

    def register(self, type_name, parser_instance):
        """Registers a parser instance for a given type name."""
        if not isinstance(parser_instance, BaseParser):
            raise TypeError("Parser instance must be a subclass of BaseParser.")
        self.parsers[type_name] = parser_instance

    def get_parser(self, type_name):
        """Retrieves a registered parser by type name."""
        parser = self.parsers.get(type_name)
        if not parser:
            raise ValueError(f"No parser registered for type '{type_name}'.")
        return parser

class PDFParser(BaseParser):
    """Parses PDF files using PyMuPDF (fitz)."""
    def parse(self, file_path, indices=None):
        """Parses the PDF file and extracts text and image information.
        `indices` can be a string (e.g., "1,3-5,N") specifying page numbers or ranges.
        """
        if fitz is None:
            raise ParsingError("PyMuPDF (fitz) is not installed. Please install it to parse PDF files. You can typically install it with: pip install PyMuPDF")

        try:
            doc = fitz.open(file_path)
            text_parts = []
            images_info = []
            
            num_pages_total = doc.page_count
            pages_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                pages_to_process_indices = parse_index_string(indices, num_pages_total)
                if not pages_to_process_indices and indices.strip(): # If string was not empty but parsing yielded no pages
                    # This could be due to invalid indices like "999" for a 10-page doc
                    # Or a malformed string like "abc"
                    # parse_index_string prints warnings for unparseable parts.
                    # If the result is empty, it means no valid pages were selected.
                    # We could choose to process all, or none, or raise error.
                    # For now, let's process no pages if specific indices were given but resulted in empty set.
                    print(f"Warning: PDF index string '{indices}' resulted in no pages to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list): # Support direct list of 0-indexed integers (internal use?)
                pages_to_process_indices = [p for p in indices if 0 <= p < num_pages_total]
            else: # No indices or unhandled type, process all pages
                pages_to_process_indices = list(range(num_pages_total))
            
            if not pages_to_process_indices and num_pages_total > 0 and indices:
                # If indices were provided, but resulted in an empty list, extract nothing.
                # (This condition is now partly handled above, but good for clarity)
                pass # text_parts and images_info will remain empty
            elif not pages_to_process_indices and num_pages_total > 0 and not indices:
                 # If no indices provided and pages exist, process all (already default by range)
                 pages_to_process_indices = list(range(num_pages_total))
            elif not pages_to_process_indices and num_pages_total == 0:
                # No pages in doc, nothing to process
                pass 

            for page_num_0_indexed in pages_to_process_indices:
                # Page numbers are already validated by parse_index_string to be within bounds
                page = doc.load_page(page_num_0_indexed)
                text_parts.append(page.get_text("text").strip())
                
                img_list = page.get_images(full=True)
                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    images_info.append({
                        "page_num": page_num_0_indexed + 1, # Store as 1-indexed for consistency in output
                        "img_index_on_page": img_index,
                        "xref": xref,
                        "width": img[2],
                        "height": img[3],
                        "format": img[5],
                        "bbox": page.get_image_bbox(img).irect.to_json()
                    })
            
            doc.close()
            text_content = "\n\n".join(text_parts)

            return {
                "text": text_content,
                "images": images_info,
                "num_pages": num_pages_total, # Total pages in doc
                "file_path": file_path,
                # Store 1-indexed pages that were actually processed
                "indices_processed": [p + 1 for p in pages_to_process_indices]
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PDF: File not found at {file_path}")
        except Exception as e:
            if "no such file or directory" in str(e).lower() or \
               "cannot open" in str(e).lower():
                raise ParsingError(f"Error parsing PDF: File not found or cannot be opened at {file_path}. (PyMuPDF: {e})")
            elif "damaged" in str(e).lower() or "cannot be opened" in str(e).lower():
                 raise ParsingError(f"Error reading PDF file {file_path}. The file might be corrupted or encrypted. (PyMuPDF: {e})")
            raise ParsingError(f"An unexpected error occurred while parsing PDF file {file_path} with PyMuPDF: {e}")

class PPTXParser(BaseParser):
    """Parses PowerPoint (PPTX) files."""
    def parse(self, file_path, indices=None):
        """Parses the PPTX file and extracts text content from selected slides.
        `indices` can be a string (e.g., "1,3-5,N") specifying slide numbers or ranges.
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            all_slides = list(prs.slides) # Convert to list to use indices
            num_slides_total = len(all_slides)
            slides_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                slides_to_process_indices = parse_index_string(indices, num_slides_total)
                if not slides_to_process_indices and indices.strip():
                    print(f"Warning: PPTX index string '{indices}' resulted in no slides to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list):
                slides_to_process_indices = [s for s in indices if 0 <= s < num_slides_total]
            else: # No indices or unhandled type, process all slides
                slides_to_process_indices = list(range(num_slides_total))

            text_content_parts = []
            processed_slide_numbers_1_indexed = []

            if not slides_to_process_indices and num_slides_total > 0 and indices:
                 pass # Extract nothing if indices given but empty result
            elif not slides_to_process_indices and num_slides_total > 0 and not indices:
                 slides_to_process_indices = list(range(num_slides_total)) # Process all if no indices
            elif not slides_to_process_indices and num_slides_total == 0:
                 pass # No slides in doc

            for slide_num_0_indexed in slides_to_process_indices:
                slide = all_slides[slide_num_0_indexed]
                slide_text_parts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text_parts.append(run.text)
                
                if slide_text_parts:
                    # Using original 1-indexed number for slide header for clarity
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n{' '.join(slide_text_parts)}")
                else: # Add a marker for blank slides if they were selected
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n[Blank Slide or No Text Content]")
                processed_slide_numbers_1_indexed.append(slide_num_0_indexed + 1)

            final_text_content = "\n\n".join(text_content_parts).strip()

            return {
                "text": final_text_content,
                "num_slides": num_slides_total, # Total slides in doc
                "file_path": file_path,
                "indices_processed": processed_slide_numbers_1_indexed # 1-indexed slides processed
            }
        except ImportError:
            raise ParsingError("python-pptx is not installed. Please install it to parse PPTX files. You can typically install it with: pip install python-pptx")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PPTX: File not found at {file_path}")
        except Exception as e:
            # Check if it's a PackageNotFoundError, which might indicate a corrupted file for python-pptx
            if "PackageNotFoundError" in str(type(e)):
                 raise ParsingError(f"Error parsing PPTX file {file_path}. The file might be corrupted, not a valid PPTX, or an issue with python-pptx: {e}")
            raise ParsingError(f"An unexpected error occurred while parsing PPTX file {file_path}: {e}")

class HTMLParser(BaseParser):
    """Parses HTML files into Markdown text using html2text."""
    def parse(self, file_path, indices=None):
        """Parses the HTML file and converts its content to Markdown.
        `indices` is currently ignored for HTML files.
        """
        try:
            import html2text
            h = html2text.HTML2Text()
            # Configure html2text options if needed, e.g.:
            # h.ignore_links = True
            # h.ignore_images = True
            # h.body_width = 0 # No line wrapping

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            markdown_content = h.handle(html_content)
            
            # For HTML, num_pages/num_slides isn't directly applicable in the same way.
            # We can omit it or add other relevant metadata if available (e.g., title).
            # For now, keep it simple.
            return {
                "text": markdown_content.strip(),
                "file_path": file_path,
                # "indices_processed": [] # Or None, as indices are not used yet
            }
        except ImportError:
            raise ParsingError("html2text is not installed. Please install it to parse HTML files. You can typically install it with: pip install html2text")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing HTML: File not found at {file_path}")
        except Exception as e:
            raise ParsingError(f"An unexpected error occurred while parsing HTML file {file_path}: {e}")

# Example of how parsers might be registered (this would typically happen in the Attachments core or user code)
# parser_registry = ParserRegistry()
# parser_registry.register('pdf', PDFParser())
# parser_registry.register('pptx', PPTXParser()) 