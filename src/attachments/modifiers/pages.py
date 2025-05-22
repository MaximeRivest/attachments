"""Page selection modifier for PDF documents and PowerPoint presentations."""

import io
from typing import List, Optional, Any

try:
    from pypdf import PdfReader, PdfWriter
    _PdfReader = PdfReader
    _PdfWriter = PdfWriter
except ImportError:
    _PdfReader = None  # type: ignore
    _PdfWriter = None  # type: ignore

try:
    from pptx.presentation import Presentation
    from pptx import Presentation as PresentationConstructor
    _Presentation = Presentation
    _PresentationConstructor = PresentationConstructor
except ImportError:
    _Presentation = None  # type: ignore
    _PresentationConstructor = None  # type: ignore

from ..core.decorators import modifier


if _PdfReader and _PdfWriter:
    @modifier
    def pages(pdf_reader: _PdfReader, page_spec: Optional[str] = None) -> _PdfReader:
        """
        Select specific pages from PDF using pypdf (MIT-compatible).
        
        Args:
            pdf_reader: pypdf PdfReader object
            page_spec: Page specification string like "1,3-5,-1" or ":3"
                      If None, uses commands from attachment
                      
        Returns:
            New PdfReader with selected pages
            
        Examples:
            pages(pdf, "1,3-5")    # Pages 1, 3, 4, 5
            pages(pdf, ":3")       # First 3 pages  
            pages(pdf, "-1")       # Last page
            pages(pdf, "N")        # Last page (alternative)
        """
        # Get page specification from attachment commands if not provided
        if page_spec is None:
            if hasattr(pdf_reader, '_attachment_commands'):
                page_spec = pdf_reader._attachment_commands.get('pages', '')
            else:
                return pdf_reader  # No page specification, return original
        
        if not page_spec:
            return pdf_reader
            
        total_pages = len(pdf_reader.pages)
        wanted_pages: List[int] = []
        
        # Parse page specification (same logic as your original transform)
        page_spec = page_spec.replace(" ", "")
        specs = page_spec.split(",")
        
        for spec in specs:
            if spec == "":  # ignore stray commas
                continue
                
            if spec.startswith(":"):  # leading range :3 == first 3 pages
                end_user_indexed = int(spec[1:]) if spec[1:] else total_pages
                end_0_indexed = min(end_user_indexed, total_pages)
                wanted_pages.extend(range(0, end_0_indexed))
                continue
                
            if spec.upper() == "N" or spec == "-1":  # Last page
                if total_pages > 0:
                    wanted_pages.append(total_pages - 1)
                continue
                
            if "-" in spec:  # Range e.g., "3-5" (pages 3, 4, 5)
                start_str, end_str = spec.split("-")
                start_0_indexed = int(start_str) - 1  # Convert to 0-indexed
                end_0_indexed = int(end_str) - 1      # Convert to 0-indexed (inclusive)
                
                if 0 <= start_0_indexed < total_pages and start_0_indexed <= end_0_indexed:
                    wanted_pages.extend(range(start_0_indexed, min(end_0_indexed + 1, total_pages)))
                continue
                
            # Single page "4" (page 4, which is index 3)
            page_0_indexed = int(spec) - 1
            if 0 <= page_0_indexed < total_pages:
                wanted_pages.append(page_0_indexed)
        
        # Remove duplicates and sort
        wanted_pages = sorted(list(set(p for p in wanted_pages if 0 <= p < total_pages)))
        
        # If no pages selected or all pages selected, return original
        if not wanted_pages or len(wanted_pages) == total_pages:
            return pdf_reader
        
        # Create new PDF with selected pages
        if _PdfWriter is None:
            raise ImportError("pypdf PdfWriter not available")
        
        writer = _PdfWriter()
        for page_index in wanted_pages:
            writer.add_page(pdf_reader.pages[page_index])
        
        # Write to in-memory stream and return new PdfReader
        output_stream = io.BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)
        
        if _PdfReader is None:
            raise ImportError("pypdf PdfReader not available")
        
        return _PdfReader(output_stream)


if _Presentation and _PresentationConstructor:
    @modifier
    def pages(presentation: _Presentation, page_spec: Optional[str] = None) -> _Presentation:
        """
        Select specific slides from PowerPoint presentation.
        
        Args:
            presentation: python-pptx Presentation object
            page_spec: Page specification string like "1,3-5,-1" or ":3"
                      If None, uses commands from attachment
                      
        Returns:
            New Presentation with selected slides
            
        Examples:
            pages(ppt, "1,3-5")    # Slides 1, 3, 4, 5
            pages(ppt, ":3")       # First 3 slides  
            pages(ppt, "-1")       # Last slide
            pages(ppt, "N")        # Last slide (alternative)
        """
        # Get page specification from attachment commands if not provided
        if page_spec is None:
            if hasattr(presentation, '_attachment_commands'):
                page_spec = presentation._attachment_commands.get('pages', '')
            else:
                return presentation  # No page specification, return original
        
        if not page_spec:
            return presentation
            
        slides = presentation.slides
        total_slides = len(slides)
        wanted_slides: List[int] = []
        
        # Parse page specification (same logic as PDF)
        page_spec = page_spec.replace(" ", "")
        specs = page_spec.split(",")
        
        for spec in specs:
            if spec == "":  # ignore stray commas
                continue
                
            if spec.startswith(":"):  # leading range :3 == first 3 slides
                end_user_indexed = int(spec[1:]) if spec[1:] else total_slides
                end_0_indexed = min(end_user_indexed, total_slides)
                wanted_slides.extend(range(0, end_0_indexed))
                continue
                
            if spec.upper() == "N" or spec == "-1":  # Last slide
                if total_slides > 0:
                    wanted_slides.append(total_slides - 1)
                continue
                
            if "-" in spec:  # Range e.g., "3-5" (slides 3, 4, 5)
                start_str, end_str = spec.split("-")
                start_0_indexed = int(start_str) - 1  # Convert to 0-indexed
                end_0_indexed = int(end_str) - 1      # Convert to 0-indexed (inclusive)
                
                if 0 <= start_0_indexed < total_slides and start_0_indexed <= end_0_indexed:
                    wanted_slides.extend(range(start_0_indexed, min(end_0_indexed + 1, total_slides)))
                continue
                
            # Single slide "4" (slide 4, which is index 3)
            slide_0_indexed = int(spec) - 1
            if 0 <= slide_0_indexed < total_slides:
                wanted_slides.append(slide_0_indexed)
        
        # Remove duplicates and sort
        wanted_slides = sorted(list(set(s for s in wanted_slides if 0 <= s < total_slides)))
        
        # If no slides selected or all slides selected, return original
        if not wanted_slides or len(wanted_slides) == total_slides:
            return presentation
        
        # Create new presentation with selected slides
        if _PresentationConstructor is None:
            raise ImportError("python-pptx Presentation not available")
        new_presentation = _PresentationConstructor()
        
        # Copy slide master from original (preserves theme, layouts)
        # Note: This is a simplified approach. In a full implementation,
        # we might need to copy the slide master more carefully
        
        # Remove the default blank slide that comes with new presentations
        # Note: This is a simplified approach. We could implement proper slide removal later
        slides_to_remove = list(new_presentation.slides)
        for slide in slides_to_remove:
            # For now, we'll work around the default slide by not removing it
            # and instead overwriting content as we add new slides
            pass
        
        # Copy selected slides
        for slide_index in wanted_slides:
            source_slide = slides[slide_index]
            
            # Get the slide layout from the source slide
            source_layout = source_slide.slide_layout
            
            # Find a matching layout in the new presentation or use the first one
            target_layout = new_presentation.slide_layouts[0]  # Simplified: use first layout
            
            # Add a new slide with this layout
            new_slide = new_presentation.slides.add_slide(target_layout)
            
            # Copy slide content
            # This is a simplified copy - in a full implementation we'd need to
            # copy all shapes, text, images, etc. more carefully
            try:
                # Try to copy the slide title if it exists
                if source_slide.shapes.title and new_slide.shapes.title:
                    new_slide.shapes.title.text = source_slide.shapes.title.text
            except:
                pass  # Continue if title copying fails
        
        return new_presentation

