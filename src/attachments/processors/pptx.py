"""Processor for PowerPoint presentations (.pptx).

Extracts per-slide text (titles as ``# Title`` headings, body text frames,
and tables) plus, on request, embedded pictures. Slide boundaries are
published as ``meta.segments`` per the IR contract.

Requires ``python-pptx``: ``pip install attachments[pptx]``.

Legacy binary ``.ppt`` files are NOT supported — python-pptx only reads
the Office Open XML (``.pptx``) format, so no processor is registered
for ``.ppt``. Slide rasterization (rendering slides to PNG) needs an
external renderer (LibreOffice/PowerPoint) and is out of scope; only
EMBEDDED pictures are extracted.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from ..options import Option, register_options
from ..types import ERROR_PARSE, error_artifact, make_artifact, missing_dep_artifact
from . import register_processor

log = logging.getLogger("attachments.processors.pptx")

_SLIDE_SEPARATOR = "\n\n"


def _join_slides_with_segments(
    slides: list[tuple[str, str]],
) -> tuple[str, list[dict]]:
    """Join per-slide texts and build slide segments (IR contract: meta.segments).

    *slides* is a list of ``(label, slide_text)`` pairs. Each segment
    carries the slide label and start/end offsets that slice the joined
    text exactly back to that slide's text.

    Examples:
        >>> text, segs = _join_slides_with_segments(
        ...     [("Intro", "# Intro\\nhello"), ("slide 2", "world")]
        ... )
        >>> text
        '# Intro\\nhello\\n\\nworld'
        >>> segs[0]
        {'kind': 'slide', 'label': 'Intro', 'start': 0, 'end': 13}
        >>> text[segs[1]["start"] : segs[1]["end"]]
        'world'
    """
    parts: list[str] = []
    segments: list[dict] = []
    offset = 0
    for label, slide_text in slides:
        if parts:
            offset += len(_SLIDE_SEPARATOR)
        segments.append(
            {
                "kind": "slide",
                "label": label,
                "start": offset,
                "end": offset + len(slide_text),
            }
        )
        parts.append(slide_text)
        offset += len(slide_text)
    return _SLIDE_SEPARATOR.join(parts), segments


def _walk_shapes(shapes):
    """Yield shapes depth-first, descending into group shapes.

    Group shapes expose a nested ``.shapes`` collection; everything else
    is yielded as-is (duck typing keeps this import-free).
    """
    for shape in shapes:
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _walk_shapes(nested)
        else:
            yield shape


def _table_text(table) -> str:
    """Render a pptx table as a simple pipe-delimited text table."""
    rows: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    return "\n".join(rows)


def _slide_text(slide) -> tuple[str, str | None]:
    """Return ``(slide_text, title)`` for one slide.

    The title (when present and non-empty) becomes a ``# <title>``
    heading line; remaining text frames and tables follow in shape
    order, joined with newlines.
    """
    try:
        title_shape = slide.shapes.title
    except Exception:  # malformed layouts must never break extraction
        title_shape = None

    title: str | None = None
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title = title_shape.text_frame.text.strip() or None

    parts: list[str] = []
    if title:
        parts.append(f"# {title}")
    for shape in _walk_shapes(slide.shapes):
        # BaseShape.__eq__ compares the underlying XML element ("is" would
        # fail: python-pptx builds fresh shape wrappers on every iteration).
        if title_shape is not None and shape == title_shape:
            continue  # already emitted as the heading line
        if getattr(shape, "has_table", False):
            parts.append(_table_text(shape.table))
        elif getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts), title


def pptx_processor(
    data: bytes,
    *,
    filename: str | None = None,
    render_images: bool | str = False,
    **_opts: Any,
) -> dict[str, Any]:
    """PowerPoint (.pptx) -> artifact dict (text, images, audio, video, meta).

    Options (via att(..., **options)):
      - render_images: truthy (True/"always"/"auto") extracts EMBEDDED
        pictures as image items with their 1-based slide number as
        ``page``. Default False. Slides are never rasterized.

    meta.kind is ``"slides"``; meta.segments has one ``"slide"`` segment
    per slide whose offsets slice ``text`` exactly to that slide's text.
    """
    source = filename or "file.pptx"

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return missing_dep_artifact(source, "pptx")

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        log.warning("failed to open pptx %s: %s", source, e)
        artifact = error_artifact(
            source, ERROR_PARSE, f"Failed to parse PowerPoint file: {e}"
        )
        artifact["meta"]["kind"] = "slides"
        return artifact

    try:
        # --- per-slide text + segments ---
        entries: list[tuple[str, str]] = []
        for slide_no, slide in enumerate(prs.slides, start=1):
            slide_text, title = _slide_text(slide)
            entries.append((title or f"slide {slide_no}", slide_text))
        text, segments = _join_slides_with_segments(entries)

        extra: dict[str, Any] = {"slide_count": len(entries)}

        # --- embedded picture extraction (optional) ---
        images: list[dict[str, Any]] = []
        if render_images:
            stem = filename or "presentation"
            for slide_no, slide in enumerate(prs.slides, start=1):
                for shape in _walk_shapes(slide.shapes):
                    try:
                        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE and not hasattr(
                            shape, "image"
                        ):
                            continue
                        # Raises for linked (not embedded) pictures.
                        image = shape.image
                        mimetype = image.content_type or "image/png"
                        ext = image.ext or mimetype.split("/")[-1].split("+")[0]
                        images.append(
                            {
                                "name": (
                                    f"{stem}-slide-{slide_no}-"
                                    f"image-{len(images) + 1}.{ext}"
                                ),
                                "mimetype": mimetype,
                                "bytes": image.blob,
                                "page": slide_no,
                            }
                        )
                    except Exception as exc:
                        log.debug("skipping picture on slide %d: %s", slide_no, exc)
            extra["images_extracted"] = len(images)
    except Exception as e:
        log.warning("failed to extract pptx content from %s: %s", source, e)
        artifact = error_artifact(
            source, ERROR_PARSE, f"Failed to parse PowerPoint file: {e}"
        )
        artifact["meta"]["kind"] = "slides"
        return artifact

    meta: dict[str, Any] = {"kind": "slides", "extra": extra}
    if segments:
        meta["segments"] = segments
    return make_artifact(text=text, images=images, meta=meta)


OPTIONS: tuple[Option, ...] = (
    Option(
        "images",
        "bool_or_auto",
        aliases=("render",),
        param="render_images",
        default=False,
        help=(
            "Extract embedded slide pictures: true/false "
            "('auto'/'always' behave like true; slides are never rasterized)."
        ),
        example="images: true",
    ),
)


def register() -> None:
    """Register the pptx processor and its option schema (idempotent)."""
    register_processor(".pptx", pptx_processor)
    register_options(".pptx", OPTIONS)


register()
