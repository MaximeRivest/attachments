"""File parsing logic."""

from abc import ABC, abstractmethod
import os
import io
import re
from typing import Any, Optional, Union

from PIL import Image, UnidentifiedImageError, ImageOps
# Removed: from pillow_heif import register_heif_opener

import fitz  # PyMuPDF for PDFParser
from pptx import Presentation # for PPTXParser
from bs4 import BeautifulSoup # for HTMLParser
import mammoth # for DOCXParser (docx to html)
from odf import text as odf_text, teletype as odf_teletype # for ODTParser
from odf.opendocument import load as load_odf # for ODTParser
from .utils import parse_index_string
from .utils import parse_audio_operations
from .exceptions import ParsingError, DetectionError # Ensure these are correctly imported from local package
from .config import Config # Added for parser config access

from typing import Any
# Conditional imports for optional audio features
try:
    from pydub import AudioSegment
    from pydub.exceptions import CouldntDecodeError
    PYDUB_AVAILABLE = True
except ImportError:
    AudioSegment = None # type: ignore
    CouldntDecodeError = None # type: ignore
    PYDUB_AVAILABLE = False

try:
    import speech_recognition as sr
    SPEECH_RECOGNITION_AVAILABLE = True
except ImportError:
    sr = None # type: ignore
    SPEECH_RECOGNITION_AVAILABLE = False

# Conditional import for pillow_heif
try:
    import pillow_heif
    HAS_PILLOW_HEIF = True
    pillow_heif.register_heif_opener() # Call only if import succeeded
except ImportError:
    pillow_heif = None # type: ignore
    HAS_PILLOW_HEIF = False

# Global constants for audio processing (can be moved to config or class if preferred)
SUPPORTED_AUDIO_FORMATS = ['wav', 'mp3', 'ogg', 'flac', 'm4a', 'aac']
DEFAULT_AUDIO_FORMAT = 'wav'
DEFAULT_AUDIO_BITRATE = "192k"

class BaseParser(ABC):
    """Abstract base class for file parsers."""
    @abstractmethod
    def parse(self, file_path, indices=None):
        """Parses the file and returns structured content."""
        pass

class ParserRegistry:
    """Manages registration and retrieval of parsers."""
    def __init__(self):
        self.parsers = {}

    def register(self, type_name, parser_instance):
        """Registers a parser instance for a given type name."""
        if not isinstance(parser_instance, BaseParser):
            raise TypeError("Parser instance must be a subclass of BaseParser.")
        self.parsers[type_name] = parser_instance

    def get_parser(self, type_name):
        """Retrieves a registered parser by type name."""
        parser = self.parsers.get(type_name)
        if not parser:
            raise ValueError(f"No parser registered for type '{type_name}'.")
        return parser

class PDFParser(BaseParser):
    """Parses PDF files using PyMuPDF (fitz)."""
    def parse(self, file_path, indices=None):
        """Parses the PDF file and extracts text and image information.
        `indices` can be a string (e.g., "1,3-5,N") specifying page numbers or ranges.
        """
        if fitz is None:
            raise ParsingError("PyMuPDF (fitz) is not installed. Please install it to parse PDF files. You can typically install it with: pip install PyMuPDF")

        try:
            doc = fitz.open(file_path)
            text_parts = []
            images_info = []
            
            num_pages_total = doc.page_count
            pages_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                pages_to_process_indices = parse_index_string(indices, num_pages_total)
                if not pages_to_process_indices and indices.strip(): # If string was not empty but parsing yielded no pages
                    # This could be due to invalid indices like "999" for a 10-page doc
                    # Or a malformed string like "abc"
                    # parse_index_string prints warnings for unparseable parts.
                    # If the result is empty, it means no valid pages were selected.
                    # We could choose to process all, or none, or raise error.
                    # For now, let's process no pages if specific indices were given but resulted in empty set.
                    print(f"Warning: PDF index string '{indices}' resulted in no pages to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list): # Support direct list of 0-indexed integers (internal use?)
                pages_to_process_indices = [p for p in indices if 0 <= p < num_pages_total]
            else: # No indices or unhandled type, process all pages
                pages_to_process_indices = list(range(num_pages_total))
            
            if not pages_to_process_indices and num_pages_total > 0 and indices:
                # If indices were provided, but resulted in an empty list, extract nothing.
                # (This condition is now partly handled above, but good for clarity)
                pass # text_parts and images_info will remain empty
            elif not pages_to_process_indices and num_pages_total > 0 and not indices:
                 # If no indices provided and pages exist, process all (already default by range)
                 pages_to_process_indices = list(range(num_pages_total))
            elif not pages_to_process_indices and num_pages_total == 0:
                # No pages in doc, nothing to process
                pass 

            for page_num_0_indexed in pages_to_process_indices:
                # Page numbers are already validated by parse_index_string to be within bounds
                page = doc.load_page(page_num_0_indexed)
                text_parts.append(page.get_text("text").strip())
                
                img_list = page.get_images(full=True)
                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    images_info.append({
                        "page_num": page_num_0_indexed + 1, # Store as 1-indexed for consistency in output
                        "img_index_on_page": img_index,
                        "xref": xref,
                        "width": img[2],
                        "height": img[3],
                        "format": img[5],
                        "bbox": page.get_image_bbox(img).irect.to_json()
                    })
            
            doc.close()
            text_content = "\n\n".join(text_parts)

            return {
                "type": "pdf",
                "text": text_content,
                "images": images_info,
                "page_count": num_pages_total, # Renamed from num_pages for consistency
                "file_path": file_path,
                "indices_processed": [p + 1 for p in pages_to_process_indices]
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PDF: File not found at {file_path}")
        except Exception as e:
            if "no such file or directory" in str(e).lower() or \
               "cannot open" in str(e).lower():
                raise ParsingError(f"Error parsing PDF: File not found or cannot be opened at {file_path}. (PyMuPDF: {e})")
            elif "damaged" in str(e).lower() or "cannot be opened" in str(e).lower():
                 raise ParsingError(f"Error reading PDF file {file_path}. The file might be corrupted or encrypted. (PyMuPDF: {e})")
            raise ParsingError(f"An unexpected error occurred while parsing PDF file {file_path} with PyMuPDF: {e}")

class PPTXParser(BaseParser):
    """Parses PowerPoint (PPTX) files."""
    def parse(self, file_path, indices=None):
        """Parses the PPTX file and extracts text content from selected slides.
        `indices` can be a string (e.g., "1,3-5,N") specifying slide numbers or ranges.
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            all_slides = list(prs.slides) # Convert to list to use indices
            num_slides_total = len(all_slides)
            slides_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                slides_to_process_indices = parse_index_string(indices, num_slides_total)
                if not slides_to_process_indices and indices.strip():
                    print(f"Warning: PPTX index string '{indices}' resulted in no slides to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list):
                slides_to_process_indices = [s for s in indices if 0 <= s < num_slides_total]
            else: # No indices or unhandled type, process all slides
                slides_to_process_indices = list(range(num_slides_total))

            text_content_parts = []
            processed_slide_numbers_1_indexed = []

            if not slides_to_process_indices and num_slides_total > 0 and indices:
                 pass # Extract nothing if indices given but empty result
            elif not slides_to_process_indices and num_slides_total > 0 and not indices:
                 slides_to_process_indices = list(range(num_slides_total)) # Process all if no indices
            elif not slides_to_process_indices and num_slides_total == 0:
                 pass # No slides in doc

            for slide_num_0_indexed in slides_to_process_indices:
                slide = all_slides[slide_num_0_indexed]
                slide_text_parts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text_parts.append(run.text)
                
                if slide_text_parts:
                    # Using original 1-indexed number for slide header for clarity
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n{' '.join(slide_text_parts)}")
                else: # Add a marker for blank slides if they were selected
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n[Blank Slide or No Text Content]")
                processed_slide_numbers_1_indexed.append(slide_num_0_indexed + 1)

            final_text_content = "\n\n".join(text_content_parts).strip()

            return {
                "type": "pptx",
                "text": final_text_content,
                "num_slides": num_slides_total, # Total slides in doc
                "file_path": file_path,
                "indices_processed": processed_slide_numbers_1_indexed # 1-indexed slides processed
            }
        except ImportError:
            raise ParsingError("python-pptx is not installed. Please install it to parse PPTX files. You can typically install it with: pip install python-pptx")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PPTX: File not found at {file_path}")
        except Exception as e:
            # Check if it's a PackageNotFoundError, which might indicate a corrupted file for python-pptx
            if "PackageNotFoundError" in str(type(e)):
                 raise ParsingError(f"Error parsing PPTX file {file_path}. The file might be corrupted, not a valid PPTX, or an issue with python-pptx: {e}")
            raise ParsingError(f"An unexpected error occurred while parsing PPTX file {file_path}: {e}")

class HTMLParser(BaseParser):
    """Parses HTML files into Markdown text using html2text."""
    def parse(self, file_path, indices=None):
        """Parses the HTML file and converts its content to Markdown.
        `indices` is currently ignored for HTML files.
        """
        try:
            import html2text
            h = html2text.HTML2Text()
            # Configure html2text options if needed, e.g.:
            # h.ignore_links = True
            # h.ignore_images = True
            # h.body_width = 0 # No line wrapping

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            markdown_content = h.handle(html_content)
            
            # For HTML, num_pages/num_slides isn't directly applicable in the same way.
            # We can omit it or add other relevant metadata if available (e.g., title).
            # For now, keep it simple.
            return {
                "type": "html",
                "text": markdown_content.strip(),
                "file_path": file_path,
                # "indices_processed": [] # Or None, as indices are not used yet
            }
        except ImportError:
            raise ParsingError("html2text is not installed. Please install it to parse HTML files. You can typically install it with: pip install html2text")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing HTML: File not found at {file_path}")
        except Exception as e:
            raise ParsingError(f"An unexpected error occurred while parsing HTML file {file_path}: {e}")

class DOCXParser(BaseParser):
    """Parses Word (DOCX) files using python-docx."""
    def parse(self, file_path, indices=None):
        """Parses the DOCX file and extracts text content.
        `indices` is currently ignored for DOCX files.
        """
        try:
            from docx import Document
        except ImportError:
            raise ParsingError("python-docx is not installed. Please install it to parse DOCX files. You can typically install it with: pip install python-docx")

        try:
            document = Document(file_path)
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            # Add text from tables if needed
            # for table in document.tables:
            #     for row in table.rows:
            #         for cell in row.cells:
            #             full_text.append(cell.text)
            text_content = '\n'.join(full_text)

            return {
                "type": "docx",
                "text": text_content.strip(),
                "file_path": file_path,
                # "indices_processed": [] # Or None, as indices are not used yet
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing DOCX: File not found at {file_path}")
        except Exception as e:
            if "Package not found" in str(e) or "doesn't exist" in str(e): # More robust check for python-docx file errors
                raise ParsingError(f"Error parsing DOCX file {file_path}. The file might be corrupted or not a valid DOCX: {e}")
            raise ParsingError(f"An unexpected error occurred while parsing DOCX file {file_path}: {e}")

class ODTParser(BaseParser):
    """Parses OpenDocument Text (ODT) files using odfpy."""
    def parse(self, file_path, indices=None):
        """Parses the ODT file and extracts text content.
        `indices` is currently ignored for ODT files.
        """
        try:
            from odf.opendocument import load
            from odf import text as odftext # Import the text module specifically
        except ImportError:
            raise ParsingError("odfpy is not installed. Please install it to parse ODT files. You can typically install it with: pip install odfpy")

        try:
            doc = load(file_path)
            all_text_elements = doc.getElementsByType(odftext.P) # Use odftext.P
            full_text = []
            for te in all_text_elements:
                # Iterate through child nodes to extract actual text content
                text_content_for_element = ""
                for child in te.childNodes:
                    if child.nodeType == child.TEXT_NODE:
                        text_content_for_element += child.data
                full_text.append(text_content_for_element.strip())
            
            final_text_content = '\n'.join(filter(None, full_text)) # Filter out empty strings from blank paragraphs

            return {
                "type": "odt",
                "text": final_text_content.strip(),
                "file_path": file_path,
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing ODT: File not found at {file_path}")
        except Exception as e:
            # odfpy might raise various errors for corrupted files, e.g. zipfile.BadZipFile
            if "BadZipFile" in str(type(e)) or "not a valid ODF file" in str(e):
                 raise ParsingError(f"Error parsing ODT file {file_path}. The file might be corrupted or not a valid ODT: {e}")
            raise ParsingError(f"An unexpected error occurred while parsing ODT file {file_path}: {e}")

class ImageParser(BaseParser):
    """Parses various image formats using Pillow, with special handling for HEIC/HEIF."""

    def __init__(self, config: Optional[Config] = None):
        self.config = config if config else Config()

    def _parse_image_operations(self, img: Image.Image, ops_str: str) -> tuple[Image.Image, dict, list[str], str | None, int | None]:
        img_modified = img.copy()
        applied_ops_summary = {}
        text_parts = []
        current_output_format = None
        current_quality = None

        operations = ops_str.lower().split(',')
        for op_full_str in operations:
            op_full_str = op_full_str.strip()
            if not op_full_str:
                continue

            parts = op_full_str.split(':', 1)
            op_name = parts[0].strip()
            op_value_str = parts[1].strip() if len(parts) > 1 else ""

            try:
                if op_name == "resize":
                    if 'x' not in op_value_str:
                        print(f"Warning: Invalid resize format '{op_value_str}'. Expected 'WIDTHxHEIGHT'. Skipping resize op.")
                        continue
                    width_str, height_str = op_value_str.split('x', 1)
                    
                    parsed_width: Optional[int] = None # Explicitly Optional[int]
                    if width_str != "auto":
                        try:
                            parsed_width = int(width_str)
                            if parsed_width <= 0: raise ValueError("Width must be positive.")
                        except ValueError:
                            print(f"Warning: Invalid width '{width_str}' for resize. Must be positive integer or 'auto'. Skipping resize op.")
                            continue
                    
                    parsed_height: Optional[int] = None # Explicitly Optional[int]
                    if height_str != "auto":
                        try:
                            parsed_height = int(height_str)
                            if parsed_height <= 0: raise ValueError("Height must be positive.")
                        except ValueError:
                            print(f"Warning: Invalid height '{height_str}' for resize. Must be positive integer or 'auto'. Skipping resize op.")
                            continue

                    if parsed_width is None and parsed_height is None:
                        print("Warning: Resize specified with auto x auto. No change. Skipping resize op.")
                        continue

                    resize_op_value_w = width_str if width_str == "auto" else parsed_width
                    resize_op_value_h = height_str if height_str == "auto" else parsed_height
                    applied_ops_summary['resize'] = (resize_op_value_w, resize_op_value_h)
                    text_parts.append(f"resize:{width_str}x{height_str}")

                    original_w, original_h = img_modified.size
                    target_w: Optional[int] = parsed_width # Start with parsed values
                    target_h: Optional[int] = parsed_height

                    if target_w is None and target_h is not None: # height is specified, width is auto
                        aspect_ratio = original_w / original_h if original_h else 1 # Avoid division by zero
                        target_w = int(target_h * aspect_ratio)
                    elif target_h is None and target_w is not None: # width is specified, height is auto
                        aspect_ratio = original_h / original_w if original_w else 1 # Avoid division by zero
                        target_h = int(target_w * aspect_ratio)
                    # If both are None, we already continued. If both are not None, they are used directly.
                    
                    # Ensure target_w and target_h are now integers before comparison/use
                    if not isinstance(target_w, int) or not isinstance(target_h, int):
                        print(f"Warning: Could not determine valid integer dimensions for resize ({target_w}x{target_h}). Skipping resize op.")
                        continue
                        
                    if target_w <= 0 or target_h <= 0:
                        print(f"Warning: Calculated resize dimensions resulted in non-positive value ({target_w}x{target_h}). Skipping resize op.")
                        continue
                        
                    img_modified = img_modified.resize((target_w, target_h), Image.Resampling.LANCZOS)
                elif op_name == "rotate":
                    try:
                        angle = float(op_value_str)
                        img_modified = img_modified.rotate(angle, expand=True)
                        applied_ops_summary['rotate'] = angle
                        text_parts.append(f"rotate:{op_value_str}")
                    except ValueError:
                        print(f"Warning: Invalid angle '{op_value_str}' for rotate. Must be a number. Skipping rotate op.")
                
                elif op_name == "format":
                    fmt = op_value_str.lower()
                    if fmt in ['jpg', 'jpeg', 'png', 'webp', 'gif', 'bmp', 'tiff']: 
                        current_output_format = 'jpeg' if fmt == 'jpg' else fmt
                        applied_ops_summary['format'] = current_output_format
                        text_parts.append(f"format:{fmt}")
                    else:
                        print(f"Warning: Unsupported format '{fmt}'. Skipping format op.")

                elif op_name == "quality":
                    try:
                        q = int(op_value_str)
                        if not (0 <= q <= 100):
                            raise ValueError("Quality must be between 0 and 100.")
                        current_quality = q
                        applied_ops_summary['quality'] = q
                        text_parts.append(f"quality:{q}")
                    except ValueError as e_qual:
                        print(f"Warning: Invalid quality '{op_value_str}'. {e_qual}. Skipping quality op.")
                
                elif op_name == "grayscale" or op_name == "greyscale":
                    img_modified = ImageOps.grayscale(img_modified)
                    applied_ops_summary['grayscale'] = True
                    text_parts.append("grayscale")

                else:
                    print(f"Warning: Unknown image operation '{op_name}'. Skipping.")
            
            except Exception as e_op:
                print(f"Error during image operation '{op_full_str}': {e_op}. Skipping this operation.")
                continue

        return img_modified, applied_ops_summary, text_parts, current_output_format, current_quality

    def parse(self, file_path, indices=None):
        parsed_data = {} 
        current_img = None
        original_format = None
        original_mode = None
        original_width, original_height = 0, 0
        
        file_ext_lower = os.path.splitext(file_path)[1].lower()

        # Explicit check for HEIC/HEIF and missing dependency first
        if file_ext_lower in ['.heic', '.heif'] and not HAS_PILLOW_HEIF:
            raise ParsingError(
                f"Processing {file_ext_lower.upper()} file '{file_path}' requires the 'pillow_heif' library. "
                f"Please install it by running: pip install pillow-heif"
            )

        try:
            try:
                img_opened_by_pillow = Image.open(file_path)
                original_format = img_opened_by_pillow.format
                original_mode = img_opened_by_pillow.mode
                original_width, original_height = img_opened_by_pillow.size
                img_opened_by_pillow.load() 
                current_img = img_opened_by_pillow.copy() 
            except UnidentifiedImageError:
                # Pillow couldn't identify it. If it was HEIC/HEIF, the check above should have caught it if pillow_heif was missing.
                # If pillow_heif IS available, let the specific HEIC block try.
                if file_ext_lower in ['.heic', '.heif'] and HAS_PILLOW_HEIF and pillow_heif:
                    pass # Will be handled by the HEIC block below
                else:
                    # For other unidentified types, or if it was HEIC and pillow_heif IS installed but still couldn't parse (unlikely here)
                    raise ParsingError(f"Cannot identify image file (unsupported format or corrupt): {file_path}")
            except FileNotFoundError:
                 raise ParsingError(f"Image file not found: {file_path}")

            # Specific handling for HEIC/HEIF if Pillow failed or if format is explicitly HEIC
            # This block now assumes HAS_PILLOW_HEIF is True if file_ext_lower suggests HEIC, due to the check at the start of the method.
            if (file_ext_lower in ['.heic', '.heif'] or (original_format and original_format.upper() in ['HEIC', 'HEIF'])):
                if not (HAS_PILLOW_HEIF and pillow_heif):
                    # This case should ideally be caught by the top-level check, but as a safeguard:
                    raise ParsingError(
                        f"Internal Error: Attempting HEIC/HEIF processing for '{file_path}' but pillow_heif is not available. "
                        f"Please install it: pip install pillow-heif"
                    )
                try:
                    heif_file = pillow_heif.read_heif(file_path)
                    img_from_heif = Image.frombytes(
                        heif_file.mode,
                        heif_file.size,
                        heif_file.data,
                        "raw",
                    )
                    if img_from_heif.mode not in ('RGB', 'RGBA', 'L'):
                        img_from_heif = img_from_heif.convert('RGBA' if 'A' in heif_file.mode else 'RGB')
                    
                    current_img = img_from_heif 
                    original_format = "HEIF" 
                    original_mode = current_img.mode 
                    original_width, original_height = current_img.size
                    parsed_data['output_format_for_base64'] = 'png' 
                except Exception as e_heif:
                    raise ParsingError(f"Error parsing HEIC/HEIF file {file_path} with pillow_heif: {e_heif}") from e_heif
            
            if current_img is None:
                # This should ideally not be reached if logic is correct, means some image type wasn't handled
                raise ParsingError(f"Could not open or identify image file (and not a recognized HEIC needing pillow_heif): {file_path}")

        except FileNotFoundError: 
            raise ParsingError(f"Image file not found: {file_path}")
        except ParsingError: # Re-raise parsing errors explicitly
            raise
        except Exception as e: # Catch other general errors during initial load
            raise ParsingError(f"Error loading image file {file_path}: {e}")

        applied_ops_details = {}
        ops_log_parts = []
        
        output_format = original_format.lower() if original_format else self.config.DEFAULT_IMAGE_OUTPUT_FORMAT
        output_quality = self.config.DEFAULT_IMAGE_QUALITY

        if indices: 
            try:
                current_img, applied_ops_details, ops_log_parts, op_output_format, op_quality = \
                    self._parse_image_operations(current_img.copy(), indices)
                if op_output_format:
                    output_format = op_output_format
                if op_quality is not None:
                    output_quality = op_quality
            except Exception as e_ops_parse:
                print(f"Warning: Could not apply one or more image operations from '{indices}' for {file_path}: {e_ops_parse}. Using image as is or with partial ops.")

        final_width, final_height = current_img.size
        
        text_summary_parts = [
            f"Image: {os.path.basename(file_path)}.",
            f"Original format: {original_format}, Original mode: {original_mode}, Original dims: {original_width}x{original_height}."
        ]
        if applied_ops_details:
            op_summary_str = ", ".join(ops_log_parts) if ops_log_parts else "details unavailable"
            if not ops_log_parts and applied_ops_details:
                temp_ops_summary = []
                for k,v in applied_ops_details.items():
                    if isinstance(v, tuple) and len(v) == 2: temp_ops_summary.append(f"{k}:{v[0]}x{v[1]}")
                    else: temp_ops_summary.append(f"{k}:{v}")
                op_summary_str = ", ".join(temp_ops_summary)

            text_summary_parts.append(f"Operations: {op_summary_str}.")
            text_summary_parts.append(f"Final dims: {final_width}x{final_height}, Output as: {output_format}.")
        else:
            text_summary_parts.append(f"Final dims: {final_width}x{final_height} (no operations).")
        
        text_summary = " ".join(text_summary_parts)

        result_type = original_format.lower() if original_format else 'image' 
        if result_type in ['heic', 'heif']: result_type = 'heif' 

        final_data = {
            "type": result_type,
            "text": text_summary,
            "image_object": current_img,
            "original_format": original_format, 
            "original_mode": original_mode,
            "original_dimensions": (original_width, original_height),
            "dimensions_after_ops": (final_width, final_height),
            "operations_applied": applied_ops_details if applied_ops_details else False, 
            "output_format": output_format,
            "output_quality": output_quality,
            "file_path": file_path,
            "original_basename": os.path.basename(file_path)
        }
        if 'output_format_for_base64' in parsed_data:
            final_data['output_format_for_base64'] = parsed_data['output_format_for_base64']
            
        return final_data

class AudioParser(BaseParser):
    """Parses audio files, applying transformations using pydub."""
    def parse(self, file_path, indices=None):
        """Parses the audio file, applies transformations based on `indices` string.
        `indices` string can contain operations like "format:FMT,samplerate:SR,channels:CH,bitrate:BR".
        Example: "format:ogg,samplerate:16000,channels:1"
        """
        if AudioSegment is None:
            raise ParsingError("pydub is not installed. Please install it to process audio. pip install pydub")

        original_filename = os.path.basename(file_path)
        processed_filename_for_api = original_filename # Default, changes if format op

        try:
            # Determine original format for pydub. AudioSegment.from_file can often guess,
            # but explicit format can be more reliable if available from detection phase.
            # For now, let pydub guess or rely on file extension.
            segment = AudioSegment.from_file(file_path)
        except FileNotFoundError:
            raise ParsingError(f"Audio file not found: {file_path}") from None # Explicitly no cause if we want to hide FileNotFoundError
        except CouldntDecodeError as e:
            raise ParsingError(f"Could not decode audio file {file_path} with pydub. It might be corrupted or an unsupported format. Error: {e}") from e
        except Exception as e:
            # This catches other potential errors from AudioSegment.from_file (e.g., various OS or ffmpeg issues not covered by CouldntDecodeError)
            raise ParsingError(f"Error loading audio {file_path} with pydub: {e}") from e
        
        original_info = {
            'samplerate': segment.frame_rate,
            'channels': segment.channels,
            'sample_width': segment.sample_width,
            'duration_ms': len(segment)
        }
        current_samplerate = original_info['samplerate']
        current_channels = original_info['channels']

        operations_summary_parts = []
        applied_operations = {}
        output_format = "wav"  # Default output format if not specified
        output_bitrate = None # Default

        if indices:
            # Use the internal _parse_audio_operations which has fallback logic
            op_details_dict, operations_summary_parts_from_internal_parse = self._parse_audio_operations(indices)

            if 'format' in op_details_dict:
                output_format = op_details_dict['format'].lower()
                # Ensure applied_operations and filename reflect the actual format used (could be fallback)
                applied_operations['format'] = output_format 
                base, _ = os.path.splitext(original_filename)
                processed_filename_for_api = f"{base}.{output_format}"
                # operations_summary_parts will be rebuilt based on actual application or can reuse from internal

            if 'samplerate' in op_details_dict and op_details_dict['samplerate'] is not None:
                target_sr = op_details_dict['samplerate']
                if target_sr != segment.frame_rate:
                    segment = segment.set_frame_rate(target_sr)
                    current_samplerate = segment.frame_rate
                    applied_operations['samplerate'] = current_samplerate
            
            if 'channels' in op_details_dict and op_details_dict['channels'] is not None:
                target_ch = op_details_dict['channels']
                if target_ch != segment.channels:
                    segment = segment.set_channels(target_ch)
                    current_channels = segment.channels
                    applied_operations['channels'] = current_channels

            if 'bitrate' in op_details_dict and op_details_dict['bitrate'] is not None:
                output_bitrate = op_details_dict['bitrate']
                applied_operations['bitrate'] = output_bitrate
            
            # Rebuild operations_summary_parts based on what was actually applied
            # or use operations_summary_parts_from_internal_parse if its format matches expectations for ops_text
            operations_summary_parts = [] # Reset and rebuild based on applied_operations
            if applied_operations.get('format'):
                 operations_summary_parts.append(f"format to {applied_operations['format']}")
            if applied_operations.get('samplerate') != original_info['samplerate'] and applied_operations.get('samplerate') is not None:
                 operations_summary_parts.append(f"samplerate to {applied_operations['samplerate']}Hz")
            elif applied_operations.get('samplerate') == original_info['samplerate'] and 'samplerate' in op_details_dict and op_details_dict['samplerate'] is not None:
                 operations_summary_parts.append(f"samplerate kept at {applied_operations['samplerate']}Hz") # Explicitly stated if requested

            if applied_operations.get('channels') != original_info['channels'] and applied_operations.get('channels') is not None:
                 operations_summary_parts.append(f"channels to {applied_operations['channels']}")
            elif applied_operations.get('channels') == original_info['channels'] and 'channels' in op_details_dict and op_details_dict['channels'] is not None:
                 operations_summary_parts.append(f"channels kept at {applied_operations['channels']}") # Explicitly stated if requested

            if applied_operations.get('bitrate'):
                 operations_summary_parts.append(f"bitrate {applied_operations['bitrate']}")

        else:
            # Default processing if no ops: convert to 16kHz, mono, wav
            if segment.frame_rate != 16000:
                segment = segment.set_frame_rate(16000)
                current_samplerate = 16000; operations_summary_parts.append(f"default samplerate to {current_samplerate}Hz")
                applied_operations['samplerate'] = current_samplerate
            if segment.channels != 1:
                segment = segment.set_channels(1)
                current_channels = 1; operations_summary_parts.append(f"default to mono ({current_channels} channel)")
                applied_operations['channels'] = current_channels
            # Default output_format is already 'wav', filename reflects this
            if not processed_filename_for_api.endswith('.wav'):
                base, _ = os.path.splitext(original_filename)
                processed_filename_for_api = f"{base}.wav"
                applied_operations['format'] = 'wav' # Ensure format is recorded if defaulted

        ops_text = f"Processed with ops: [{', '.join(operations_summary_parts)}]." if operations_summary_parts else "Processed with default settings."
        text_summary = (
            f"[Audio: {original_filename}] "
            f"Original: {original_info['samplerate']}Hz, {original_info['channels']}ch, {original_info['duration_ms']}ms. "
            f"{ops_text} "
            f"Output as: {output_format}, {current_samplerate}Hz, {current_channels}ch."
        )

        return {
            "type": "audio",
            "text": text_summary,
            "file_path": file_path,
            "original_basename": original_filename, # Needed for .audios property consistency
            "audio_segment": segment, # The pydub AudioSegment object, transformed
            "output_format": output_format,
            "output_samplerate": current_samplerate,
            "output_channels": current_channels,
            "output_bitrate": output_bitrate, # e.g., "128k"
            "processed_filename_for_api": processed_filename_for_api,
            "applied_operations": applied_operations
        } 

    def _parse_audio_operations(self, ops_str: str) -> tuple[dict, list[str]]:
        op_details = {
            'format': DEFAULT_AUDIO_FORMAT, # Default format
            'samplerate': None, 
            'channels': None,
            'bitrate': DEFAULT_AUDIO_BITRATE # Default bitrate
        }
        text_parts: list[str] = []
        applied_ops: dict[str, Any] = {}

        if not ops_str: # No operations string provided
            return op_details, text_parts

        operations = ops_str.split(',')
        for op_detail in operations:
            parts = op_detail.split(':', 1)
            if len(parts) != 2:
                # print(f"Warning: Invalid operation format '{op_detail}'. Skipping.")
                continue
            op_key = parts[0].strip().lower()
            op_value = parts[1].strip()

            if op_key == 'format':
                requested_format = op_value.lower()
                if requested_format in SUPPORTED_AUDIO_FORMATS:
                    op_details['format'] = requested_format
                    text_parts.append(f"format to {requested_format}")
                    applied_ops[op_key] = requested_format
                else:
                    # print(f"Warning: Unsupported audio format '{requested_format}'. Defaulting to '{DEFAULT_AUDIO_FORMAT}'.")
                    op_details['format'] = DEFAULT_AUDIO_FORMAT # Fallback
                    text_parts.append(f"format to {DEFAULT_AUDIO_FORMAT} (requested '{requested_format}' unsupported)")
                    applied_ops[op_key] = DEFAULT_AUDIO_FORMAT # Store the actual format used
            elif op_key == 'samplerate':
                try:
                    target_sr = int(op_value)
                    op_details['samplerate'] = target_sr
                    text_parts.append(f"samplerate to {target_sr}Hz")
                    applied_ops[op_key] = target_sr
                except ValueError:
                    # print(f"Warning: Invalid samplerate value '{op_value}'. Ignoring.")
                    pass
            elif op_key == 'channels':
                try:
                    target_ch = int(op_value)
                    op_details['channels'] = target_ch
                    text_parts.append(f"channels to {target_ch}")
                    applied_ops[op_key] = target_ch
                except ValueError:
                    # print(f"Warning: Invalid channel value '{op_value}'. Ignoring.")
                    pass
            elif op_key == 'bitrate':
                try:
                    output_bitrate = op_value
                    op_details['bitrate'] = output_bitrate
                    text_parts.append(f"bitrate {output_bitrate}")
                    applied_ops[op_key] = output_bitrate
                except ValueError:
                    # print(f"Warning: Invalid bitrate value '{op_value}'. Ignoring.")
                    pass

        return op_details, text_parts